#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_gamification_slides():
    """
    PDFから抽出した内容を基にGamefication講座のプレゼンテーションを作成
    """
    prs = Presentation()
    
    # スライド1: タイトルスライド
    slide_layout = prs.slide_layouts[0]  # タイトルスライド
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Why Most Gamification Fails"
    subtitle.text = "Section 1: ゲーミフィケーションが失敗する理由\nMVPバージョン (0:00-7:00)"
    
    # フォント設定
    title_font = title.text_frame.paragraphs[0].font
    title_font.size = Pt(44)
    title_font.bold = True
    title_font.color.rgb = RGBColor(51, 51, 51)
    
    subtitle_font = subtitle.text_frame.paragraphs[0].font
    subtitle_font.size = Pt(20)
    subtitle_font.color.rgb = RGBColor(102, 102, 102)
    
    # スライド2: フック - 共感できる失敗例
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツ
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "フック：共感できる失敗例 (0:00-1:30)"
    content.text = """ナレーター：

「1月1日。あなたはフィットネスアプリをダウンロードします。歩数でポイント、連続利用で
バッジ、友達とのランキングがあるアプリです。2月15日...そのアプリはスマホの3ページ目に
埋もれています。開き義気がありますか？」

[ビジュアル: 削除されたアプリ、放置された企業研修モジュール]

「あるいは、企業研修を受けているかもしれません。『このモジュールを完了して、50ポイントを獲得し
ましょう！次のレベルをアンロックしましょう！』と表示され、チェックボックスにチェックを入れるだ
けだけにクリックする。本当の学習も、永続的な変化もありません。ただ...順守しているだけなのです。

「不快な真実があります。ほとんどのゲーミフィケーションは、単に失敗するだけでなく、
状況を悪化させるのです。人間を、私が「ポイントゾンビ」と呼ぶ存在に変えてしまうので
す。そして今日は、その理由を理解し、全く異なる道を見つけるのです。」"""
    
    # スライド3: 問題分析 - なぜ失敗するのか
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "問題分析：なぜ失敗するのか (1:30-4:00)"
    content.text = """「問題はゲーミフィケーション自体ではありません。問題は、ほとんどのデザイナーが
人間の本質を根本的に誤解していることです。彼らは人々を報酬を求める機械として見
ているのではなく、創造的で意味を生み出す存在として見ているのです。」

[ビジュアル: バブロフの犬の実験 vs チェスをする子供たち]

「これは、研究者イアン・ボゴストが『ポイント化』と呼ぶ現象につながります──活動
の深い目的を理解せずに、ポイント、バッジ、リーダーボードを貼り付けることです。そ
れは、スプレッドシートにゲーム衣装を着せて魔法を期待するようなものです。」

[ビジュアル: ポイント化の一例（シンプルなポイントシステム、意味のないバッジ）]

「しかし、実際はこうです: 人間のモチベーションを外部報酬に還元すると、活動を持続
可能にする内在的な喜びを破壊してしまいます。あなたはゲーミフィケーションをしてい
るのではなく、操作しているだけです。」"""
    
    # スライド4: 解決策のプレビュー
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "解決策のプレビュー：意味のあるゲーミフィケーション (4:00-6:30)"
    content.text = """「ここで登場するのが、学習たちが『意味のあるゲーミフィケーション』と呼ぶものです。
これは『どうすれば人々を私たちの望む行動に導けるか』ではなく、『どうすれば人々が本
当に価値のあるものを発見できるよう支援できるか』を問うものです。」

[ビジュアル: 意味のあるゲーミフィケーションの例: オープンエンドの学習プラッ
トフォーム、創造的なチャレンジ]

「人間を『制約の創造的な解決者』として扱い、報酬を求めるロボットとして扱いません。
チェスのルールのような構造を提供し、創造的な表現を制限するのではなく、実際に可能
にします。」

[ビジュアル: チェスゲーム、音楽のスケール、建築の設計図など、制約が創造性
を引き出す例]

「この違いは根本的です。ポイント化は『どのように行動を制御するか』を問います、意味
のあるゲーミフィケーションは『どのように人間の主体性を高めるか』を問います。前者は
依存を生み出します。後者は自律性を生み出します。」"""
    
    # スライド5: コースロードマップ
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "コースロードマップ：学習者が発見すること (6:30-7:00)"
    content.text = """「このコースでは、操作と解放の違いを認識する方法を学びます。意味のあるゲーミフィ
ケーションの5つの原則を発見します。そして最も重要なのは、

創造性、自律性、人間性を高める体験を設計するためのツールを獲得します。

なぜなら、目標は人々をゲーム化することではないからです。目標は、彼らの中に眠る
プレイヤーを発見する手助けをすることです。

始めましょう。"

[ビジュアル: コースロゴ、「Begin Your Journey」のコールトゥアクション]"""
    
    # スライド6: 制作メモ
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "**Production Notes**"
    content.text = """**ビジュアルスタイルの推奨事項**

• トーン: プロフェッショナルだが親しみやすく、学術的だがアクセスしやすい
• グラフィック: 派手な効果ではなく、クリーンでミニマリストなアニメーション
• カラーパレット: 人間中心（温かい青、緑）vs 機械的（冷たいグレー、赤）
• タイポグラフィ: 信頼性と知性を伝える、クリアで読みやすいフォント

**必要な主要なビジュアル要素**

1. 画面分割比較：ポイント化 vs 意味のあるゲーミフィケーション
2. Real app screenshots: Examples of both approaches (with permission/fair use)
3. シンプルなアニメーション：外部動機から内在的動機への移行を示すデータ可視化
4. 技能可能なパターンと持続不可能なパターンのエンゲージメント曲線 人間の顔：本物
5. のエンゲージメント vs 強制的な順守を示す"""
    
    # プレゼンテーションを保存
    output_path = "/Users/sakaihironori/Desktop/AI_PJ/Gamification講座/Why_Most_Gamification_Fails_Section1.pptx"
    prs.save(output_path)
    print(f"プレゼンテーションが作成されました: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_gamification_slides()